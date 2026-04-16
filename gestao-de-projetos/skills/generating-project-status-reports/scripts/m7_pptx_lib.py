#!/usr/bin/env python3
"""
M7 Investimentos — PPTX Building Blocks Library (v4 — M7-2026 Brand)

A library of composable building blocks for creating M7-branded presentations.
Claude writes a Python script importing this library and composing slides creatively.

KEY DESIGN PRINCIPLE:
  - Cover, Agenda, Section Divider, Closing → DARK background (#424135 Verde Caqui)
  - All content slides → LIGHT background (#fffdef Off-White), dark text
  - Font: TWK Everett Light (fallback Arial) — headings are NOT bold
  - Accent: Lime #eef77c (decorative only, never on text)

Brand source: M7 Brandbook 2026 + multi7.com.br

Usage:
    import sys
    sys.path.insert(0, "<path-to-skill>/scripts")
    from m7_pptx_lib import M7Presentation, C

    prs = M7Presentation()
    prs.add_cover("PROJETO", "Título Principal", tagline="Frase de impacto")
    slide = prs.add_content_slide("01", "CONTEXTO", "Captação superou meta em 23%")
    slide.add_problem_items([...])
    prs.save("output.pptx")

Requires: pip install python-pptx
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

# === PATHS ===
SKILL_DIR = Path(__file__).resolve().parent.parent
ASSETS_DIR = SKILL_DIR / "assets"
LOGO_DARK_PATH = ASSETS_DIR / "m7-logo-dark.png"      # Verde caqui logo → light bg
LOGO_OFFWHITE_PATH = ASSETS_DIR / "m7-logo-offwhite.png"  # Off-white logo → dark bg
LOGO_PATH = LOGO_DARK_PATH  # Default (backward compat)
HERO_DARK_PATH = ASSETS_DIR / "m7-hero-dark.png"       # Darkened hero photo → cover bg

# === FONT ===
# Brandbook specifies TWK Everett Light. python-pptx can't embed fonts,
# so we reference the font name (requires system install) with Arial fallback.
FONT_HEADING = "twkEverett"   # Headings — Light weight, NOT bold
FONT_BODY = "twkEverett"      # Body text — Light weight
FONT_FALLBACK = "Arial"       # Universal fallback

# === COLORS (M7-2026 Brand — Brandbook + multi7.com.br) ===
class C:
    """M7-2026 Color Palette — from brandbook and multi7.com.br."""
    # Backgrounds
    DARK = RGBColor(0x42, 0x41, 0x35)           # Verde Caqui — cover/closing bg
    OFF_WHITE = RGBColor(0xFF, 0xFD, 0xEF)      # Off-White (warm) — content bg
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)           # Pure white — cards
    CARD = RGBColor(0xF6, 0xF6, 0xF5)           # Card/box background (verde-caqui-50)
    SEPARATOR = RGBColor(0xD0, 0xD0, 0xCC)      # Divider lines (verde-caqui-100)

    # Text
    TEXT_TITLE = RGBColor(0x42, 0x41, 0x35)      # Titles on light bg (= DARK)
    TEXT_BODY = RGBColor(0x4F, 0x4E, 0x3C)       # Body text (verde-medio)
    TEXT_MED = RGBColor(0x79, 0x75, 0x5C)        # Intro paragraphs (verde-claro)
    TEXT_LIGHT = RGBColor(0xAE, 0xAD, 0xA8)      # Descriptions, footnotes (verde-caqui-200)

    # Accents
    LIME = RGBColor(0xEE, 0xF7, 0x7C)           # Primary accent, section labels (DECORATIVE ONLY)
    GREEN = RGBColor(0x4C, 0xAF, 0x50)          # Positive/success/done
    AMBER = RGBColor(0xF5, 0x9E, 0x0B)          # Warning/in-progress
    RED = RGBColor(0xE4, 0x69, 0x62)            # Negative/risk/problem
    BLUE = RGBColor(0x3B, 0x82, 0xF6)           # Info/sprint accent
    BLUE_OFFICE = RGBColor(0x00, 0x70, 0xC0)    # Office blue (Webdings checkmarks)
    GREEN_DONE = RGBColor(0x00, 0xB0, 0x50)     # Milestone done (green)
    RED_DARK = RGBColor(0xC0, 0x00, 0x00)       # Dark red (overdue/late)
    PURPLE = RGBColor(0x8B, 0x5C, 0xF6)         # Tertiary accent
    TEAL = RGBColor(0x14, 0xB8, 0xA6)           # Quaternary accent

    # Callout backgrounds
    BG_HIGHLIGHT = RGBColor(0xFF, 0xFD, 0xF4)   # Highlight callout (off-white-300)
    BG_GREEN = RGBColor(0xE8, 0xF5, 0xE9)       # Light green (success)
    BG_RED = RGBColor(0xFD, 0xED, 0xED)         # Light red (warning)

    # Legacy aliases (backward compatibility)
    BG_YELLOW = RGBColor(0xFF, 0xFD, 0xF4)      # Alias for BG_HIGHLIGHT

# Accent color sequence for sequential items (steps, sprints, etc.)
ACCENT_SEQ = [C.DARK, C.LIME, C.AMBER, C.GREEN, C.BLUE, C.PURPLE, C.TEAL]

# === DIMENSIONS (from reference PPTX, in EMU) ===
SLIDE_W = 9144000    # 10 inches
SLIDE_H = 5143500    # 5.625 inches
LOGO_W = 502920      # ~0.55 inches (width)
LOGO_H = 251460      # ~0.275 inches (height) — 2:1 ratio matching actual logo assets
LOGO_SIZE = LOGO_W   # Deprecated alias (backward compat)
LOGO_LEFT = 8321040
LOGO_TOP = 137160
MARGIN = 548640      # ~0.6 inches
CONTENT_W = 8046720  # ~8.8 inches (full width minus margins)
WIDE_W = 7772400     # Width for intro text
HEADER_TOP = 365760  # Section label Y
TITLE_TOP = 640080   # Title Y
INTRO_TOP = 1371600  # Intro paragraph Y
CONTENT_TOP = 1371600  # Default content start Y (below title)
FOOTNOTE_TOP = 4617720  # Footnote Y
FOOTNOTE_W = 7772400

# === EXECUTIVE STATUS LAYOUT (compact, Votorantim-inspired) ===
EXEC_HEADER_TOP = 228600       # Section label Y (compact)
EXEC_TITLE_TOP = 365760        # Title Y (14pt)
EXEC_CONTENT_TOP = 594360      # Content start (was 1371600)
EXEC_GAP = 27432               # Tight gap between zones (matches vBruno spacing)

# ── Executive Status Static Layout (vBruno reference coordinates) ──
# All zones at FIXED EMU positions — only content varies between reports.
# These values were extracted from the manually adjusted vBruno PPTX.
_EX_LEFT = 320040              # 0.350" — wider than MARGIN, matches vBruno
_EX_WIDTH = 8503920            # 9.300" — total content width
_EX_COL_W = 4183380            # 4.575" — each column width
_EX_COL_GAP = 137160           # 0.150" — gap between columns
_EX_RIGHT_X = _EX_LEFT + _EX_COL_W + _EX_COL_GAP  # 4640580
_EX_PAD_X = 137160             # Horizontal padding inside timeline box
_EX_LABEL_OFFSET = 54864       # Label pill offset from box left edge

# Vertical positions (Y) — fixed for all Executive Status slides
_EX_TL_LABEL_Y = 1114864      # Timeline fieldset label
_EX_TL_BOX_Y = 1207298        # Timeline container box
_EX_TL_BOX_H = 722376         # Timeline box height
_EX_TL_LINE_REL_Y = 210000    # Timeline line Y relative to box top
_EX_TL_LINE_H = 13716         # Timeline line thickness

_EX_COL_LABEL_Y = 2002681     # Two-column fieldset labels
_EX_COL_BOX_Y = 2095115       # Two-column boxes top (label_y + 92434)
_EX_COL_BOX_H = 1755648       # Two-column box height

_EX_ATT_LABEL_Y = 3921761     # Attention fieldset label
_EX_ATT_BOX_Y = 4014195       # Attention box top (label_y + 92434)
_EX_ATT_BOX_H = 795528        # Attention box height

_EX_FOOT_Y = 4880000          # Footnote Y
_EX_LEGEND_Y = 4917788        # Legend Y

# Timeline diamond sizes
_EX_BIG_DOT = 164592          # Done/current diamond (~0.18")
_EX_SMALL_DOT = 118872        # Pending diamond (~0.13")
_EX_LABEL_H = 307777          # Milestone label height (allows 2-line wrap)
_EX_MIN_LABEL_W = 1295400     # Min label width (~1.42", matches vBruno)


# ──────────────────────────────────────────────
# Low-level drawing primitives
# ──────────────────────────────────────────────

def _set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _set_hero_bg(slide):
    """Set hero-dark image as full-slide background. Falls back to solid DARK."""
    if HERO_DARK_PATH.exists():
        slide.shapes.add_picture(
            str(HERO_DARK_PATH), 0, 0, SLIDE_W, SLIDE_H
        )
    else:
        _set_bg(slide, C.DARK)


def _add_logo(slide, dark_bg=False):
    """Add M7 logo. Uses off-white logo on dark bg, dark logo on light bg."""
    logo = LOGO_OFFWHITE_PATH if dark_bg else LOGO_DARK_PATH
    if not logo.exists():
        # Fallback to any available logo
        for p in [LOGO_DARK_PATH, LOGO_OFFWHITE_PATH, ASSETS_DIR / "m7-logo.png"]:
            if p.exists():
                logo = p
                break
        else:
            return
    slide.shapes.add_picture(
        str(logo), LOGO_LEFT, LOGO_TOP, LOGO_W, LOGO_H
    )


def _textbox(slide, left, top, width, height, text,
             font=FONT_BODY, size=Pt(9), bold=False, color=C.TEXT_BODY,
             align=None, wrap=True):
    """Add a text box with a single run. Returns the shape."""
    box = slide.shapes.add_textbox(left, top, width, height)
    box.fill.background()
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _rect(slide, left, top, width, height, fill_color):
    """Add a filled rectangle (no border)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _rounded_rect(slide, left, top, width, height, fill_color):
    """Add a filled rounded rectangle (no border)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _dashed_rect(slide, left, top, width, height,
                 border_color=C.SEPARATOR, border_width=Pt(0.75),
                 fill_color=None):
    """Rectangle with dashed border (Votorantim fieldset style)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    shape.line.dash_style = MSO_LINE.DASH
    return shape


def _diamond(slide, left, top, size, fill_color):
    """Small diamond shape for timeline milestones."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _circle(slide, left, top, size, fill_color, text, text_color=C.WHITE,
            font_size=Pt(9)):
    """Add a filled circle with centered text."""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill_color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(text)
    run.font.name = FONT_BODY
    run.font.size = font_size
    run.font.bold = True
    run.font.color.rgb = text_color
    return circle


def _separator(slide, left, top, width, color=C.SEPARATOR):
    _rect(slide, left, top, width, 4572, color)


# ──────────────────────────────────────────────
# Executive Status — Static layout helpers
# (Fixed-position elements, not tied to SlideBuilder cursor)
# ──────────────────────────────────────────────

def _exec_fieldset_label(slide, left, top, title, color=C.TEXT_BODY):
    """Fieldset-legend pill label at a fixed position."""
    label_w = min(len(title) * 54864 + 91440, _EX_WIDTH // 2)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left + _EX_LABEL_OFFSET, top,
        label_w, 180000
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C.WHITE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(18288)
    tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_BODY
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = color
    return shape


def _exec_content_box(slide, left, top, width, height, body_lines=None):
    """Content box with solid thin border and bold bullet text at a fixed position."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C.WHITE
    shape.line.color.rgb = C.SEPARATOR
    shape.line.width = Pt(0.75)

    if body_lines:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(91440)
        tf.margin_right = Emu(91440)
        tf.margin_top = Emu(54864)
        tf.margin_bottom = Emu(54864)
        for j, line in enumerate(body_lines):
            text = f"• {line}" if not line.startswith("•") else line
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(4) if j > 0 else Pt(0)
            p.space_after = Pt(2)
            run = p.add_run()
            run.text = text
            run.font.name = FONT_BODY
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = C.TEXT_BODY
    return shape


def _exec_timeline_markers(slide, milestones):
    """Place milestone diamonds and labels at fixed positions inside the timeline box."""
    n = len(milestones)
    if n == 0:
        return

    line_left = _EX_LEFT + _EX_PAD_X
    line_right = _EX_LEFT + _EX_WIDTH - _EX_PAD_X
    line_width = line_right - line_left
    line_top = _EX_TL_BOX_Y + _EX_TL_LINE_REL_Y + _EX_BIG_DOT // 2

    # Horizontal timeline bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, line_left, line_top, line_width, _EX_TL_LINE_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C.BLUE
    bar.line.fill.background()

    # Milestone positions
    if n == 1:
        positions = [line_left + line_width // 2]
    else:
        spacing = line_width // (n - 1)
        positions = [line_left + i * spacing for i in range(n)]

    status_colors = {
        "done": C.GREEN_DONE, "current": C.BLUE, "pending": C.SEPARATOR,
    }

    for i, milestone in enumerate(milestones):
        label = milestone[0]
        status = milestone[1] if len(milestone) > 1 else "pending"
        color = status_colors.get(status, C.SEPARATOR)

        cx = positions[i]
        dot_sz = _EX_BIG_DOT if status != "pending" else _EX_SMALL_DOT
        mx = cx - dot_sz // 2
        my = _EX_TL_BOX_Y + _EX_TL_LINE_REL_Y + _EX_BIG_DOT // 2 - dot_sz // 2

        if status == "pending":
            sh = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, mx, my, dot_sz, dot_sz)
            sh.fill.solid()
            sh.fill.fore_color.rgb = C.WHITE
            sh.line.color.rgb = C.SEPARATOR
            sh.line.width = Pt(1)
        else:
            _diamond(slide, mx, my, dot_sz, color)

        # Label below diamond
        col_w = spacing if n > 1 else line_width
        lbl_w = max(col_w, _EX_MIN_LABEL_W) if n > 1 else line_width
        lbl_left = cx - lbl_w // 2
        lbl_left = max(lbl_left, _EX_LEFT + 18288)
        lbl_right = min(lbl_left + lbl_w, _EX_LEFT + _EX_WIDTH - 18288)
        lbl_w = lbl_right - lbl_left

        lbl_top = _EX_TL_BOX_Y + _EX_TL_LINE_REL_Y + _EX_BIG_DOT + 27432
        _textbox(slide, lbl_left, lbl_top, lbl_w, _EX_LABEL_H,
                 label, size=Pt(7), color=C.TEXT_LIGHT,
                 align=PP_ALIGN.CENTER)


def _exec_legend(slide, legend_y):
    """Place the 4-item colored diamond legend at a fixed Y position."""
    legend_items = [
        ("Não iniciado", C.SEPARATOR),
        ("Em andamento", C.BLUE),
        ("Atrasado", C.RED_DARK),
        ("Concluído", C.GREEN_DONE),
    ]
    legend_w = len(legend_items) * 822960
    legend_left = (SLIDE_W - legend_w) // 2
    diamond_sz = 109728
    for i, (lbl, clr) in enumerate(legend_items):
        item_x = legend_left + i * 822960
        _diamond(slide, item_x, legend_y + 36576, diamond_sz, clr)
        _textbox(slide, item_x + diamond_sz + 27432, legend_y,
                 640080, 200736,
                 lbl, size=Pt(7), color=C.TEXT_BODY)


# ──────────────────────────────────────────────
# SlideBuilder — composable content on LIGHT slides
# ──────────────────────────────────────────────

class SlideBuilder:
    """Builder for composing content on a light-background content slide.

    All text defaults are tuned for light backgrounds (dark text).
    The `y` property tracks vertical position for auto-layout.
    """

    def __init__(self, slide, y=None):
        self._slide = slide
        self._y = y or CONTENT_TOP

    @property
    def slide(self):
        return self._slide

    @property
    def y(self):
        return self._y

    @y.setter
    def y(self, val):
        self._y = val

    # ── Intro Paragraph ──

    def add_intro(self, text, left=MARGIN, width=WIDE_W):
        """Gray explanatory paragraph below the title (TWK Everett 12pt, verde-claro)."""
        h = 731520
        _textbox(self._slide, left, self._y, width, h,
                 text, size=Pt(12), color=C.TEXT_MED)
        self._y += h
        return self

    # ── Footnote ──

    def add_footnote(self, text):
        """Gray footnote at the bottom of the slide (TWK Everett 9pt, verde-caqui-200).
        Does NOT advance y — always pinned to FOOTNOTE_TOP."""
        _textbox(self._slide, MARGIN, FOOTNOTE_TOP, FOOTNOTE_W, 274320,
                 text, size=Pt(9), color=C.TEXT_LIGHT)
        return self

    # ── Problem Items (slide 3 pattern) ──

    def add_problem_items(self, items):
        """List of problems with red accent bar + × icon.

        items: list of (title, description) tuples.
        Pattern: red left bar | × | bold title | gray description
        """
        ROW_H = 548640
        GAP = 0
        BAR_W = 54864
        ICON_LEFT = MARGIN + BAR_W + 91440
        TEXT_LEFT = ICON_LEFT + 365760

        for title, desc in items:
            # Red accent bar
            _rect(self._slide, MARGIN, self._y, BAR_W, 457200, C.RED)
            # × icon
            _textbox(self._slide, ICON_LEFT, self._y, 274320, 457200,
                     "×", size=Pt(14), bold=True, color=C.RED)
            # Title
            _textbox(self._slide, TEXT_LEFT, self._y, 2743200, 228600,
                     title, size=Pt(11), bold=True, color=C.TEXT_BODY)
            # Description
            _textbox(self._slide, TEXT_LEFT, self._y + 201168, 7315200, 228600,
                     desc, size=Pt(10), color=C.TEXT_LIGHT)
            self._y += ROW_H + GAP

        return self

    # ── Step Cards (slide 4 pattern — horizontal cards with top-bar) ──

    def add_step_cards(self, steps, colors=None):
        """Horizontal cards with colored top accent bars and arrows.

        steps: list of (number_label, title) tuples.
        colors: optional list of accent colors per card.
        """
        n = len(steps)
        if n == 0:
            return self
        if colors is None:
            colors = [ACCENT_SEQ[i % len(ACCENT_SEQ)] for i in range(n)]

        gap = 109728
        card_w = (CONTENT_W - (n - 1) * gap) // n
        card_h = 914400
        bar_h = 54864

        for i, (num_label, title) in enumerate(steps):
            x = MARGIN + i * (card_w + gap)
            color = colors[i]

            # Card background
            _rect(self._slide, x, self._y, card_w, card_h, C.CARD)
            # Top accent bar
            _rect(self._slide, x, self._y, card_w, bar_h, color)
            # Number/label
            _textbox(self._slide, x, self._y + bar_h + 68580, card_w, 365760,
                     num_label, font=FONT_HEADING, size=Pt(18), bold=True,
                     color=color, align=PP_ALIGN.CENTER)
            # Title
            _textbox(self._slide, x, self._y + bar_h + 457200, card_w, 365760,
                     title, size=Pt(10), bold=True, color=C.TEXT_BODY,
                     align=PP_ALIGN.CENTER)

            # Arrow between cards
            if i < n - 1:
                arrow_x = x + card_w
                _textbox(self._slide, arrow_x, self._y + card_h // 3,
                         gap, 228600, "→",
                         size=Pt(14), color=C.TEXT_LIGHT,
                         align=PP_ALIGN.CENTER)

        self._y += card_h + 91440
        return self

    # ── Metric Cards (slide 4 bottom pattern — big number + label) ──

    def add_metric_cards(self, metrics, colors=None):
        """Row of metric cards with left accent bar, big number, and description.

        metrics: list of (value, description) tuples.
        colors: optional list of accent colors per card.
        """
        n = len(metrics)
        if n == 0:
            return self
        if colors is None:
            colors = [C.LIME] * n

        gap = 91440
        card_w = (CONTENT_W - (n - 1) * gap) // n
        card_h = 822960
        bar_w = 54864

        for i, (value, desc) in enumerate(metrics):
            x = MARGIN + i * (card_w + gap)
            color = colors[i]

            # Card background
            _rect(self._slide, x, self._y, card_w, card_h, C.CARD)
            # Left accent bar
            _rect(self._slide, x, self._y, bar_w, card_h, color)
            # Big number
            _textbox(self._slide, x + bar_w + 91440, self._y + 137160,
                     card_w - bar_w - 182880, 365760,
                     value, font=FONT_HEADING, size=Pt(20), bold=True,
                     color=C.TEXT_TITLE)
            # Description
            _textbox(self._slide, x + bar_w + 91440, self._y + 502920,
                     card_w - bar_w - 182880, 274320,
                     desc, size=Pt(10), color=C.TEXT_LIGHT)

        self._y += card_h + 91440
        return self

    # ── Numbered Steps (slide 5 pattern — circles + title + desc + right label) ──

    def add_numbered_steps(self, steps, colors=None):
        """Vertical numbered steps with colored circles, title, description, right label.

        steps: list of (title, description, right_label) tuples.
        colors: optional list of colors per step circle.
        """
        n = len(steps)
        if n == 0:
            return self
        if colors is None:
            colors = [ACCENT_SEQ[i % len(ACCENT_SEQ)] for i in range(n)]

        circle_size = 365760
        row_h = 502920

        for i, step_data in enumerate(steps):
            title = step_data[0]
            desc = step_data[1] if len(step_data) > 1 else ""
            right_label = step_data[2] if len(step_data) > 2 else ""
            color = colors[i]

            # Numbered circle
            _circle(self._slide, MARGIN, self._y, circle_size, color, str(i + 1),
                    font_size=Pt(14))
            # Title
            _textbox(self._slide, MARGIN + circle_size + 182880, self._y + 22860,
                     2743200, 228600,
                     title, size=Pt(13), bold=True, color=C.TEXT_BODY)
            # Description
            if desc:
                _textbox(self._slide, MARGIN + circle_size + 182880,
                         self._y + 228600, 5486400, 228600,
                         desc, size=Pt(10), color=C.TEXT_LIGHT)
            # Right label
            if right_label:
                _textbox(self._slide, 7498080, self._y + 68580,
                         1097280, 274320,
                         right_label, size=Pt(10), bold=True, color=color)

            self._y += row_h

        return self

    # ── Sprint Cards (slide 8 pattern — tall roadmap cards) ──

    def add_sprint_cards(self, sprints, colors=None):
        """Horizontal sprint/roadmap cards with top accent, code, title, date.

        sprints: list of (code, name, detail, date_range) tuples.
        colors: optional list of accent colors per card.
        """
        n = len(sprints)
        if n == 0:
            return self
        if colors is None:
            colors = [ACCENT_SEQ[i % len(ACCENT_SEQ)] for i in range(n)]

        gap = 109728
        card_w = (CONTENT_W - (n - 1) * gap) // n
        card_h = 2011680
        bar_h = 54864

        for i, (code, name, detail, date_range) in enumerate(sprints):
            x = MARGIN + i * (card_w + gap)
            color = colors[i]

            # Card background
            _rect(self._slide, x, self._y, card_w, card_h, C.CARD)
            # Top accent bar
            _rect(self._slide, x, self._y, card_w, bar_h, color)
            # Sprint code
            _textbox(self._slide, x, self._y + bar_h + 91440,
                     card_w, 274320,
                     code, font=FONT_HEADING, size=Pt(12), bold=True,
                     color=color, align=PP_ALIGN.CENTER)
            # Sprint name
            _textbox(self._slide, x, self._y + bar_h + 411480,
                     card_w, 228600,
                     name, size=Pt(11), bold=True, color=C.TEXT_BODY,
                     align=PP_ALIGN.CENTER)
            # Detail
            if detail:
                _textbox(self._slide, x, self._y + bar_h + 685800,
                         card_w, 365760,
                         detail, size=Pt(9), color=C.TEXT_LIGHT,
                         align=PP_ALIGN.CENTER)
            # Divider
            div_y = self._y + card_h - 365760
            _rect(self._slide, x + 182880, div_y, card_w - 365760, 9144, C.SEPARATOR)
            # Date range
            _textbox(self._slide, x, div_y + 45720,
                     card_w, 274320,
                     date_range, size=Pt(10), bold=True, color=color,
                     align=PP_ALIGN.CENTER)

            # Arrow between cards
            if i < n - 1:
                arrow_x = x + card_w
                _textbox(self._slide, arrow_x, self._y + card_h // 3,
                         gap, 228600, "→",
                         size=Pt(12), color=C.TEXT_LIGHT,
                         align=PP_ALIGN.CENTER)

        self._y += card_h + 91440
        return self

    # ── Flow Columns (slide 6 pattern — colored headers + bullets) ──

    def add_flow_columns(self, columns, colors=None):
        """Columns with colored headers, arrow connectors, and bullet items.

        columns: list of (header_text, [bullet_items]) tuples.
        colors: optional list of header background colors.
        """
        n = len(columns)
        if n == 0:
            return self
        if colors is None:
            colors = [ACCENT_SEQ[i % len(ACCENT_SEQ)] for i in range(n)]

        gap = 137160
        arrow_w = 137160
        usable_w = CONTENT_W - (n - 1) * (gap + arrow_w)
        col_w = usable_w // n
        header_h = 640080
        bullet_spacing = 347472

        # Calculate max bullets for height
        max_bullets = max(len(c[1]) for c in columns)
        total_h = header_h + max_bullets * bullet_spacing + 91440

        for i, (header, bullets) in enumerate(columns):
            x = MARGIN + i * (col_w + gap + arrow_w)
            color = colors[i]

            # Colored header
            header_shape = _rect(self._slide, x, self._y, col_w, header_h, color)
            tf = header_shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(91440)
            tf.margin_top = Emu(91440)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = header
            run.font.name = FONT_BODY
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = C.WHITE

            # Bullets below
            for j, bullet in enumerate(bullets):
                bullet_y = self._y + header_h + 91440 + j * bullet_spacing
                _textbox(self._slide, x + 45720, bullet_y,
                         col_w - 91440, 274320,
                         f"▸ {bullet}", size=Pt(10), color=C.TEXT_BODY)

            # Arrow between columns
            if i < n - 1:
                arrow_x = x + col_w + gap // 2
                _textbox(self._slide, arrow_x, self._y + header_h // 3,
                         arrow_w + gap, 365760, "→",
                         size=Pt(20), bold=True, color=C.TEXT_LIGHT,
                         align=PP_ALIGN.CENTER)

        self._y += total_h
        return self

    # ── Hierarchy Levels (slide 7 pattern — cascading N1→N4) ──

    def add_hierarchy_levels(self, levels):
        """Cascading hierarchy levels (N1 → N2 → N3 → N4).

        levels: list of (badge_label, title, description, example_text) tuples.
        Colors auto-assigned from ACCENT_SEQ.
        """
        badge_w = 457200
        badge_h = 594360
        bar_h = badge_h
        indent_step = 274320
        row_gap = 91440

        colors = [ACCENT_SEQ[i % len(ACCENT_SEQ)] for i in range(len(levels))]

        for i, level_data in enumerate(levels):
            badge_label = level_data[0]
            title = level_data[1]
            desc = level_data[2] if len(level_data) > 2 else ""
            example = level_data[3] if len(level_data) > 3 else ""
            color = colors[i]
            indent = MARGIN + i * indent_step

            # Badge (Nx)
            badge = _rect(self._slide, indent, self._y, badge_w, badge_h, color)
            tf = badge.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = badge_label
            run.font.name = FONT_HEADING
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = C.WHITE

            # Content bar
            bar_left = indent + badge_w + 9144
            bar_w = CONTENT_W - (indent - MARGIN) - badge_w - 9144
            _rect(self._slide, bar_left, self._y, bar_w, bar_h, C.CARD)

            # Title
            _textbox(self._slide, bar_left + 91440, self._y + 68580,
                     bar_w - 274320, 274320,
                     title, size=Pt(12), bold=True, color=C.TEXT_BODY)
            # Description
            if desc:
                _textbox(self._slide, bar_left + 91440, self._y + 320040,
                         bar_w - 274320, 228600,
                         desc, size=Pt(9), color=C.TEXT_LIGHT)
            # Example (right-aligned, colored)
            if example:
                _textbox(self._slide, bar_left + bar_w - 2286000,
                         self._y + 68580, 2194560, 274320,
                         example, size=Pt(10), color=color)

            # Connector dot
            if i < len(levels) - 1:
                dot_x = indent + badge_w // 2 + indent_step // 2
                _rect(self._slide, dot_x, self._y + badge_h,
                      27432, row_gap, C.SEPARATOR)

            self._y += badge_h + row_gap

        return self

    # ── Callout Boxes (slide 5/6/8 pattern — tinted bg + left accent bar) ──

    def add_callout(self, text, bg=C.BG_YELLOW, accent=C.LIME,
                    text_color=C.TEXT_BODY):
        """Callout box with tinted background, colored left bar, bold text."""
        h = 457200
        bar_w = 54864
        # Background
        _rect(self._slide, MARGIN, self._y, WIDE_W, h, bg)
        # Left accent bar
        _rect(self._slide, MARGIN, self._y, bar_w, h, accent)
        # Text
        _textbox(self._slide, MARGIN + bar_w + 91440,
                 self._y + 91440, WIDE_W - bar_w - 182880, h - 182880,
                 text, size=Pt(10), bold=True, color=text_color)
        self._y += h + 91440
        return self

    def add_success_callout(self, text, sub_text=None):
        """Green success callout (light green bg, green accent, green text)."""
        h = 548640 if sub_text else 457200
        bar_w = 54864
        _rect(self._slide, MARGIN, self._y, WIDE_W, h, C.BG_GREEN)
        _rect(self._slide, MARGIN, self._y, bar_w, h, C.GREEN)
        _textbox(self._slide, MARGIN + bar_w + 91440,
                 self._y + 91440, WIDE_W - bar_w - 182880, 274320,
                 text, size=Pt(11), bold=True, color=C.GREEN)
        if sub_text:
            _textbox(self._slide, MARGIN + bar_w + 91440,
                     self._y + 365760, WIDE_W - bar_w - 182880, 182880,
                     sub_text, size=Pt(10), color=C.TEXT_LIGHT)
        self._y += h + 91440
        return self

    def add_warning_callout(self, text):
        """Yellow warning callout."""
        return self.add_callout(text, bg=C.BG_YELLOW, accent=C.AMBER,
                                text_color=C.TEXT_BODY)

    # ── Cards (generic, on light bg) ──

    def add_card(self, title, body_lines=None, left=MARGIN, width=CONTENT_W,
                 height=None, bg=C.CARD):
        """Card with title and optional bullet lines."""
        n_lines = len(body_lines) if body_lines else 0
        if height is None:
            height = 320040 + n_lines * 200000
            height = min(height, 2800000)

        card = _rounded_rect(self._slide, left, self._y, width, height, bg)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(91440)
        tf.margin_right = Emu(91440)
        tf.margin_top = Emu(68580)
        tf.margin_bottom = Emu(68580)

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = C.TEXT_BODY

        if body_lines:
            for line in body_lines:
                p = tf.add_paragraph()
                p.space_before = Pt(3)
                run = p.add_run()
                run.text = f"• {line}" if not line.startswith("•") else line
                run.font.name = FONT_BODY
                run.font.size = Pt(9)
                run.font.color.rgb = C.TEXT_BODY

        self._y += height + 91440
        return card

    def add_status_card(self, title, body_lines=None, status="neutral",
                        left=MARGIN, width=CONTENT_W, height=None):
        """Card with colored left border. status: positive/negative/warning/neutral"""
        colors = {
            "positive": C.GREEN, "negative": C.RED,
            "warning": C.AMBER, "neutral": C.LIME,
        }
        accent = colors.get(status, C.LIME)
        bar_w = 54864
        n_lines = len(body_lines) if body_lines else 0
        if height is None:
            height = 320040 + n_lines * 200000
            height = min(height, 2800000)
        _rect(self._slide, left, self._y, bar_w, height, accent)
        return self.add_card(title, body_lines, left=left + bar_w + 9144,
                             width=width - bar_w - 9144, height=height)

    # ── Two/Three Column Layouts ──

    def add_two_columns(self, left_title, left_items, right_title, right_items,
                        left_bg=C.CARD, right_bg=C.WHITE):
        """Two side-by-side cards."""
        col_w = (CONTENT_W - 182880) // 2
        h = max(320040 + len(left_items) * 200000,
                320040 + len(right_items) * 200000)
        h = min(h, 3200000)
        saved = self._y
        self.add_card(left_title, left_items, left=MARGIN, width=col_w, height=h, bg=left_bg)
        self._y = saved
        self.add_card(right_title, right_items, left=MARGIN + col_w + 182880,
                      width=col_w, height=h, bg=right_bg)

    def add_three_columns(self, columns, bg=C.CARD):
        """Three side-by-side cards. columns: list of (title, items) tuples."""
        gap = 137160
        col_w = (CONTENT_W - 2 * gap) // 3
        max_items = max(len(c[1]) for c in columns) if columns else 0
        h = min(320040 + max_items * 200000, 3200000)
        saved = self._y
        for i, (title, items) in enumerate(columns):
            self._y = saved
            x = MARGIN + i * (col_w + gap)
            self.add_card(title, items, left=x, width=col_w, height=h, bg=bg)
        self._y = saved + h + 91440

    # ── Executive Sections (Votorantim Report Quinzenal style) ──
    #
    # Based on analysis of reference PPTX:
    #   - Section labels: AUTO_SHAPE rect with text (fieldset legend)
    #   - Content boxes: SOLID thin border (1pt), no fill
    #   - Timeline line: SQUARE_DOT dashed, 1.5pt
    #   - Content text: bold, justified, ~10pt

    def _fieldset_label(self, left, title, color=C.TEXT_BODY):
        """Fieldset-legend label that sits above a content box border."""
        # Matches vBruno reference: rounded rect with solid white fill
        label_w = min(len(title) * 54864 + 91440, CONTENT_W // 2)
        shape = self._slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left + 54864, self._y,
            label_w, 182880  # 0.20" height (matches vBruno)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = C.WHITE
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = False
        tf.margin_left = Emu(18288)
        tf.margin_top = Emu(0)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = FONT_BODY
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = color
        return label_w

    def _content_box(self, left, width, height, body_lines=None):
        """Content box with solid thin border and bold bullet text.

        Matches vBruno: AUTO_SHAPE, fill=solid white, line=0.75pt solid,
        text frame with bullet items in bold.
        """
        shape = self._slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, self._y, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = C.WHITE
        shape.line.color.rgb = C.SEPARATOR
        shape.line.width = Pt(0.75)

        if body_lines:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(91440)
            tf.margin_right = Emu(91440)
            tf.margin_top = Emu(54864)
            tf.margin_bottom = Emu(54864)
            for j, line in enumerate(body_lines):
                text = f"• {line}" if not line.startswith("•") else line
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.space_before = Pt(4) if j > 0 else Pt(0)
                p.space_after = Pt(2)
                run = p.add_run()
                run.text = text
                run.font.name = FONT_BODY
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.color.rgb = C.TEXT_BODY

        return shape

    def add_exec_section(self, title, body_lines=None, left=MARGIN,
                         width=CONTENT_W, height=None,
                         title_color=C.TEXT_BODY):
        """Fieldset-legend section with solid-border content box.

        Matches Votorantim reference: label above a thin-bordered box.
        """
        n_lines = len(body_lines) if body_lines else 0
        label_h = 164592   # fieldset label height
        gap = 27432         # gap between label and box
        if height is None:
            height = label_h + gap + max(n_lines * 228600 + 109728, 457200)
            height = min(height, 2800000)

        box_top_offset = label_h + gap
        box_h = height - box_top_offset

        # Label
        self._fieldset_label(left, title, color=title_color)
        # Content box
        saved_y = self._y
        self._y += box_top_offset
        self._content_box(left, width, box_h, body_lines)
        self._y = saved_y + height + EXEC_GAP
        return self

    def add_exec_two_columns(self, left_title, left_items,
                              right_title, right_items,
                              title_color=C.TEXT_BODY,
                              left=None, col_w=None):
        """Two side-by-side fieldset sections (Votorantim layout).

        vBruno proportions: 4.58" columns, narrow gap between them.
        Optional left/col_w override for fine positioning (defaults match vBruno).
        """
        left = left if left is not None else MARGIN
        col_w = col_w if col_w is not None else 4187952  # ~4.58"
        total_w = MARGIN + CONTENT_W - left  # available width from left edge
        gap = total_w - 2 * col_w  # remaining space as gap
        n_left = len(left_items) if left_items else 0
        n_right = len(right_items) if right_items else 0
        max_items = max(n_left, n_right, 1)
        label_h = 164592
        gap_lb = 27432
        box_h = max(max_items * 274320 + 109728, 914400)
        box_h = min(box_h, 2800000)
        total_h = label_h + gap_lb + box_h

        saved_y = self._y

        # Left column
        self._fieldset_label(left, left_title, color=title_color)
        self._y = saved_y + label_h + gap_lb
        self._content_box(left, col_w, box_h, left_items)

        # Right column
        self._y = saved_y
        right_left = left + col_w + gap
        self._fieldset_label(right_left, right_title, color=title_color)
        self._y = saved_y + label_h + gap_lb
        self._content_box(right_left, col_w, box_h, right_items)

        self._y = saved_y + total_h + EXEC_GAP
        return self

    def add_timeline(self, milestones, title="Cronograma Macro",
                     left=MARGIN, width=CONTENT_W):
        """Horizontal timeline with milestone diamonds inside a fieldset box.

        milestones: list of (label, status) tuples.
            status: "done"=filled blue, "current"=filled lime, "pending"=hollow
        Reference: solid-border box, SQUARE_DOT dashed line 1.5pt,
        diamond markers (0.23" done/current, 0.13" pending), labels below.
        """
        n = len(milestones)
        if n == 0:
            return self

        label_h = 164592       # fieldset label
        gap_lb = 27432
        pad_x = 137160         # horizontal padding inside box
        big_dot = 164592       # done/current diamond (~0.18")
        small_dot = 118872     # pending diamond (~0.13")
        date_h = 164592        # date/label row above line
        line_y_rel = date_h + 45720  # line Y relative to box top
        label_area_h = 274320  # label area below dots (allows wrapping)
        box_h = line_y_rel + big_dot + 27432 + label_area_h + 45720

        # Fieldset label
        self._fieldset_label(left, title)

        # Content box with solid border (white fill, 0.75pt — matches vBruno)
        box_top = self._y + label_h + gap_lb
        shape = self._slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, box_top, width, box_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = C.WHITE
        shape.line.color.rgb = C.SEPARATOR
        shape.line.width = Pt(0.75)

        # Timeline dashed line (SQUARE_DOT, 1.5pt — matches reference)
        line_left = left + pad_x
        line_right = left + width - pad_x
        line_width = line_right - line_left
        line_top = box_top + line_y_rel + big_dot // 2

        from pptx.oxml.ns import qn
        connector = self._slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, line_left, line_top, line_width, 13716
        )
        connector.fill.solid()
        connector.fill.fore_color.rgb = C.BLUE
        connector.line.fill.background()

        # Milestone markers + labels
        if n == 1:
            positions = [line_left + line_width // 2]
            col_w = line_width
        else:
            spacing = line_width // (n - 1)
            positions = [line_left + i * spacing for i in range(n)]
            col_w = spacing

        status_colors = {
            "done": C.GREEN_DONE, "current": C.BLUE, "pending": C.SEPARATOR,
        }

        for i, milestone in enumerate(milestones):
            label = milestone[0]
            status = milestone[1] if len(milestone) > 1 else "pending"
            color = status_colors.get(status, C.SEPARATOR)

            cx = positions[i]
            dot_sz = big_dot if status != "pending" else small_dot
            mx = cx - dot_sz // 2
            my = box_top + line_y_rel + big_dot // 2 - dot_sz // 2

            if status == "pending":
                sh = self._slide.shapes.add_shape(
                    MSO_SHAPE.DIAMOND, mx, my, dot_sz, dot_sz
                )
                sh.fill.solid()
                sh.fill.fore_color.rgb = C.WHITE
                sh.line.color.rgb = C.SEPARATOR
                sh.line.width = Pt(1)
            else:
                _diamond(self._slide, mx, my, dot_sz, color)

            # Label below diamond — centered, wrapping (min 1.42" for readability)
            lbl_w = max(col_w, 1295400) if n > 1 else line_width
            lbl_left = cx - lbl_w // 2
            lbl_left = max(lbl_left, left + 18288)
            lbl_right = min(lbl_left + lbl_w, left + width - 18288)
            lbl_w = lbl_right - lbl_left

            _textbox(self._slide, lbl_left,
                     box_top + line_y_rel + big_dot + 27432,
                     lbl_w, label_area_h,
                     label, size=Pt(7), color=C.TEXT_LIGHT,
                     align=PP_ALIGN.CENTER)

        self._y += label_h + gap_lb + box_h + EXEC_GAP
        return self

    def add_exec_attention(self, title, body_lines=None,
                            title_color=C.AMBER):
        """Attention/warning fieldset section (full width, amber label)."""
        return self.add_exec_section(
            title, body_lines,
            left=MARGIN, width=CONTENT_W,
            title_color=title_color,
        )

    # ── Metric Row (simple centered metrics) ──

    def add_metric_row(self, metrics, left=MARGIN, width=CONTENT_W):
        """Row of big centered metrics.
        metrics: list of (label, value) or (label, value, color) tuples.
        """
        n = len(metrics)
        if n == 0:
            return self
        card_w = (width - (n - 1) * 91440) // n
        card_h = 640080

        bg = _rounded_rect(self._slide, left, self._y, width, card_h, C.CARD)

        x = left + 91440
        for metric in metrics:
            label, value = metric[0], metric[1]
            val_color = metric[2] if len(metric) > 2 else C.TEXT_TITLE
            _textbox(self._slide, x, self._y + 68580, card_w, 182880,
                     label.upper(), size=Pt(7), color=C.TEXT_LIGHT,
                     align=PP_ALIGN.CENTER)
            _textbox(self._slide, x, self._y + 228600, card_w, 365760,
                     value, font=FONT_HEADING, size=Pt(20), bold=True,
                     color=val_color, align=PP_ALIGN.CENTER)
            x += card_w + 91440

        self._y += card_h + 91440
        return self

    # ── Big Number (BAN) ──

    def add_big_number(self, label, value, subtitle=None, color=C.LIME):
        """Single prominent metric. On light bg, subtitle uses TEXT_MED."""
        h = 914400 if subtitle else 731520
        _textbox(self._slide, MARGIN, self._y, CONTENT_W, 182880,
                 label.upper(), size=Pt(8), color=C.TEXT_LIGHT,
                 align=PP_ALIGN.CENTER)
        _textbox(self._slide, MARGIN, self._y + 182880, CONTENT_W, 457200,
                 value, font=FONT_HEADING, size=Pt(36), bold=True,
                 color=color, align=PP_ALIGN.CENTER)
        if subtitle:
            _textbox(self._slide, MARGIN, self._y + 594360, CONTENT_W, 228600,
                     subtitle, size=Pt(11), color=C.TEXT_MED,
                     align=PP_ALIGN.CENTER)
        self._y += h + 91440
        return self

    # ── Table ──

    def add_simple_table(self, headers, rows, left=MARGIN, width=CONTENT_W):
        """Styled table with dark header row."""
        n_rows = len(rows) + 1
        n_cols = len(headers)
        row_h = Emu(274320)
        table_h = row_h * n_rows
        shape = self._slide.shapes.add_table(
            n_rows, n_cols, left, self._y, width, table_h
        )
        table = shape.table
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = C.DARK
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(8)
                    run.font.bold = True
                    run.font.color.rgb = C.WHITE
        for i, row in enumerate(rows):
            bg = C.WHITE if i % 2 == 0 else C.CARD
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.name = FONT_BODY
                        run.font.size = Pt(8)
                        run.font.color.rgb = C.TEXT_BODY
        self._y += table_h + 91440
        return shape

    # ── Milestone Grid Table (Roadmap slide — vBruno pattern) ──

    def add_milestone_table(self, headers, rows, left=MARGIN, width=CONTENT_W):
        """Styled roadmap table with milestone columns using Webdings checkmarks.

        headers: list of column header strings (e.g., ["Sprint", "Título", "Período", "Diagnóstico", ...])
        rows: list of row dicts/lists. Each row element can be:
            - A plain string (displayed as text, Arial 8pt)
            - A tuple ("check", color) where color is an RGBColor — renders Webdings "n" (■) in that color
            - A tuple ("dash",) or just "-" — renders a centered dash

        The first 3 columns (Sprint, Título, Período) use fixed widths,
        remaining columns share the space equally for milestone checkmarks.
        """
        n_rows = len(rows) + 1
        n_cols = len(headers)
        row_h = Emu(429768)   # ~0.47" per data row (taller than default)
        header_h = Emu(274320)
        table_h = header_h + row_h * len(rows)

        # Column widths: first 3 fixed, rest equal
        fixed_widths = [Emu(530352), Emu(850392), Emu(1088136)]  # ~0.58", 0.93", 1.19"
        n_fixed = min(len(fixed_widths), n_cols)
        fixed_total = sum(fixed_widths[:n_fixed])
        remaining_w = width - fixed_total
        n_dynamic = max(n_cols - n_fixed, 1)
        dyn_w = remaining_w // n_dynamic

        col_widths = list(fixed_widths[:n_fixed]) + [Emu(dyn_w)] * n_dynamic

        shape = self._slide.shapes.add_table(
            n_rows, n_cols, left, self._y, width, table_h
        )
        table = shape.table

        # Set column widths
        for j in range(n_cols):
            table.columns[j].width = col_widths[j] if j < len(col_widths) else Emu(dyn_w)

        # Header row
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = C.DARK
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j >= n_fixed else None
                for run in p.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(8)
                    run.font.bold = True
                    run.font.color.rgb = C.WHITE

        # Data rows
        for i, row in enumerate(rows):
            bg = C.WHITE if i % 2 == 0 else C.CARD
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg

                if isinstance(val, tuple) and len(val) >= 2 and val[0] == "check":
                    # Webdings "n" = filled square checkmark
                    cell.text = ""
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = "n"
                    run.font.name = "Webdings"
                    run.font.size = Pt(16)
                    run.font.color.rgb = val[1]
                elif isinstance(val, str) and val.strip() == "-":
                    # Dash for N/A
                    cell.text = ""
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = "-"
                    run.font.name = FONT_BODY
                    run.font.size = Pt(8)
                    run.font.color.rgb = C.TEXT_BODY
                elif isinstance(val, tuple) and val[0] == "dash":
                    cell.text = ""
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = "-"
                    run.font.name = FONT_BODY
                    run.font.size = Pt(8)
                    run.font.color.rgb = C.TEXT_BODY
                else:
                    cell.text = str(val)
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.name = FONT_BODY
                            run.font.size = Pt(8)
                            run.font.color.rgb = C.TEXT_BODY

        self._y += table_h + 91440
        return shape

    # ── Highlight Box ──

    def add_highlight_box(self, text, bg=C.DARK, text_color=C.WHITE,
                          accent_word=None, accent_color=C.LIME):
        """Prominent highlight box with dark bg and optional accent word."""
        h = 365760
        box = _rounded_rect(self._slide, MARGIN, self._y, CONTENT_W, h, bg)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(137160)
        tf.margin_right = Emu(137160)
        tf.margin_top = Emu(68580)
        tf.margin_bottom = Emu(68580)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if accent_word and accent_word in text:
            parts = text.split(accent_word, 1)
            for k, part in enumerate(parts):
                if part:
                    run = p.add_run()
                    run.text = part
                    run.font.name = FONT_BODY
                    run.font.size = Pt(10)
                    run.font.color.rgb = text_color
                if k == 0:
                    run = p.add_run()
                    run.text = accent_word
                    run.font.name = FONT_BODY
                    run.font.size = Pt(10)
                    run.font.bold = True
                    run.font.color.rgb = accent_color
        else:
            run = p.add_run()
            run.text = text
            run.font.name = FONT_BODY
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = text_color
        self._y += h + 91440
        return box

    # ── Process Steps (horizontal circles) ──

    def add_process_steps(self, steps, colors=None):
        """Horizontal circles with title + description below.
        steps: list of (title, description) tuples.
        """
        n = len(steps)
        if colors is None:
            colors = [ACCENT_SEQ[i % len(ACCENT_SEQ)] for i in range(n)]
        step_w = CONTENT_W // n

        for i, (title, desc) in enumerate(steps):
            x = MARGIN + i * step_w
            cx = x + step_w // 2
            circle_size = 365760
            _circle(self._slide, cx - circle_size // 2, self._y,
                    circle_size, colors[i], str(i + 1), font_size=Pt(14))
            _textbox(self._slide, x, self._y + circle_size + 45720,
                     step_w, 228600, title,
                     size=Pt(9), bold=True, color=C.TEXT_BODY,
                     align=PP_ALIGN.CENTER)
            _textbox(self._slide, x, self._y + circle_size + 274320,
                     step_w, 365760, desc,
                     size=Pt(8), color=C.TEXT_LIGHT, align=PP_ALIGN.CENTER)
            if i < n - 1:
                _textbox(self._slide, x + step_w - 91440,
                         self._y + circle_size // 4,
                         182880, 228600, "→",
                         size=Pt(16), bold=True, color=C.TEXT_LIGHT,
                         align=PP_ALIGN.CENTER)

        self._y += 914400 + 91440
        return self

    # ── Freeform helpers ──

    def add_text(self, text, left=MARGIN, width=CONTENT_W,
                 font=FONT_BODY, size=Pt(11), bold=False, color=C.TEXT_BODY,
                 align=None):
        """Freeform text block."""
        h = 320040
        _textbox(self._slide, left, self._y, width, h,
                 text, font=font, size=size, bold=bold, color=color, align=align)
        self._y += h + 45720
        return self

    def add_spacer(self, height=91440):
        """Vertical spacing."""
        self._y += height
        return self


# ──────────────────────────────────────────────
# M7Presentation — main entry point
# ──────────────────────────────────────────────

class M7Presentation:
    """M7-2026 branded presentation builder.

    Two slide modes:
      - DARK slides: cover, agenda, section divider, closing  (bg #424135 Verde Caqui)
      - LIGHT slides: all content slides                      (bg #fffdef Off-White)
    Font: TWK Everett (Light weight for headings, not Bold).
    """

    def __init__(self):
        self._prs = Presentation()
        self._prs.slide_width = SLIDE_W
        self._prs.slide_height = SLIDE_H

    @property
    def prs(self):
        """Access the underlying python-pptx Presentation object."""
        return self._prs

    def _dark_slide(self):
        """Create a new slide with dark background (Verde Caqui #424135)."""
        slide = self._prs.slides.add_slide(self._prs.slide_layouts[6])
        _set_bg(slide, C.DARK)
        _add_logo(slide, dark_bg=True)
        return slide

    def _light_slide(self):
        """Create a new slide with light background (Off-White #fffdef)."""
        slide = self._prs.slides.add_slide(self._prs.slide_layouts[6])
        _set_bg(slide, C.OFF_WHITE)
        _add_logo(slide, dark_bg=False)
        return slide

    # ── Cover (DARK) ──

    def add_cover(self, project_label="APRESENTAÇÃO", title="Título Principal",
                  tagline=None, footer="M7 Investimentos"):
        """Cover/title slide with hero-dark image background."""
        slide = self._prs.slides.add_slide(self._prs.slide_layouts[6])
        _set_hero_bg(slide)
        _add_logo(slide, dark_bg=True)
        _textbox(slide, 0, 1463040, SLIDE_W, 365760,
                 project_label, size=Pt(10), bold=True, color=C.OFF_WHITE,
                 align=PP_ALIGN.CENTER)
        _textbox(slide, 914400, 1828800, 7315200, 1463040,
                 title, font=FONT_HEADING, size=Pt(36), bold=False,
                 color=C.OFF_WHITE, align=PP_ALIGN.CENTER)
        if tagline:
            _textbox(slide, 1371600, 3108960, 6400800, 457200,
                     tagline, size=Pt(16), color=C.LIME,
                     align=PP_ALIGN.CENTER)
        _separator(slide, 3840480, 3749040, 1463040, C.LIME)
        _textbox(slide, 1828800, 4023360, 5486400, 365760,
                 footer, size=Pt(12), color=C.OFF_WHITE,
                 align=PP_ALIGN.CENTER)
        return slide

    # ── Agenda (DARK) ──

    def add_agenda(self, items, project_label="APRESENTAÇÃO", title="Agenda"):
        """Agenda/TOC slide with dark background."""
        slide = self._dark_slide()
        _textbox(slide, MARGIN, HEADER_TOP, 4572000, 274320,
                 project_label, size=Pt(9), bold=True, color=C.OFF_WHITE)
        _textbox(slide, MARGIN, TITLE_TOP, 3657600, 548640,
                 title, font=FONT_HEADING, size=Pt(26), bold=False,
                 color=C.OFF_WHITE)
        # Numbered items with circles
        step = 384048
        y_start = 1399032
        for i, item in enumerate(items):
            y = y_start + i * step
            _circle(slide, 731520, y, 274320, C.LIME, str(i + 1))
            _textbox(slide, 1188720, y - 27432, 5486400, 320040,
                     item, size=Pt(14), color=C.OFF_WHITE)
            if i < len(items) - 1:
                _separator(slide, 1188720, y + 320040, 6858000)
        return slide

    # ── Content Slide (LIGHT) ──

    def add_content_slide(self, section_num, section_label, title, intro=None,
                          title_size=24):
        """Content slide with LIGHT background.

        Header: section label in lime, title in dark, optional intro paragraph.
        title_size: font size for title in pt (default 24, use 20 for denser slides).
        Returns SlideBuilder for composing content.
        """
        slide = self._light_slide()

        # Section label (lime accent on light bg)
        _textbox(slide, MARGIN, HEADER_TOP, 5486400, 274320,
                 f"{section_num} — {section_label}",
                 size=Pt(9), bold=True, color=C.LIME)

        # Title height adapts to font size
        title_h = 548640 if title_size >= 24 else 400110
        _textbox(slide, MARGIN, TITLE_TOP, CONTENT_W, title_h,
                 title, font=FONT_HEADING, size=Pt(title_size), bold=False,
                 color=C.TEXT_TITLE)

        # Content Y: gap after title provides breathing room
        content_y = TITLE_TOP + title_h + 457200  # ~0.50" gap after title

        sb = SlideBuilder(slide)

        # Optional intro paragraph
        if intro:
            sb.y = content_y
            sb.add_intro(intro)
        else:
            sb.y = content_y

        return sb

    # ── Section Divider (DARK) ──

    def add_section_divider(self, section_num, section_label, title):
        """Section divider (dark bg, centered)."""
        slide = self._dark_slide()
        _textbox(slide, MARGIN, 1828800, CONTENT_W, 365760,
                 f"{section_num} — {section_label}",
                 size=Pt(10), bold=True, color=C.LIME,
                 align=PP_ALIGN.CENTER)
        _textbox(slide, 914400, 2286000, 7315200, 914400,
                 title, font=FONT_HEADING, size=Pt(30), bold=False,
                 color=C.OFF_WHITE, align=PP_ALIGN.CENTER)
        return slide

    # ── Executive Status (LIGHT, static layout from vBruno) ──

    def add_executive_status_slide(self, section_num, section_label, title,
                                    milestones=None,
                                    done_items=None, next_items=None,
                                    risk_items=None, footnote=None):
        """Executive Status slide with STATIC layout (vBruno reference).

        All structural elements (boxes, labels, line, legend) are placed at
        fixed EMU coordinates (_EX_* constants). Only text content varies.

        4-zone structure:
          1. Timeline (Cronograma Macro with milestones)
          2. Status Executivo (left column)
          3. Próximas Atividades (right column)
          4. Pontos de Atenção (full-width, amber label)

        Parameters:
            section_num: e.g., "S0"
            section_label: e.g., "FUNDAÇÃO"
            title: short pattern "N de M tarefas concluídas (X%)"
            milestones: list of (label, status) tuples (max 7)
            done_items: list of strings for left column
            next_items: list of strings for right column
            risk_items: list of strings for attention zone
            footnote: string for bottom footnote
        Returns the slide object.
        """
        slide = self._light_slide()

        # ── Header (standard positions) ──
        _textbox(slide, MARGIN, HEADER_TOP, 5486400, 228600,
                 f"{section_num} — {section_label}",
                 size=Pt(9), bold=True, color=C.LIME)
        _textbox(slide, MARGIN, TITLE_TOP, CONTENT_W, 457200,
                 title, font=FONT_HEADING, size=Pt(24), bold=False,
                 color=C.TEXT_TITLE)

        # ── Zone 1: Timeline (fixed position) ──
        _exec_fieldset_label(slide, _EX_LEFT, _EX_TL_LABEL_Y,
                             "Cronograma Macro")
        _exec_content_box(slide, _EX_LEFT, _EX_TL_BOX_Y,
                          _EX_WIDTH, _EX_TL_BOX_H)
        if milestones:
            _exec_timeline_markers(slide, milestones)

        # ── Zone 2+3: Two columns (fixed position) ──
        _exec_fieldset_label(slide, _EX_LEFT, _EX_COL_LABEL_Y,
                             "Status Executivo")
        _exec_fieldset_label(slide, _EX_RIGHT_X, _EX_COL_LABEL_Y,
                             "Próximas Atividades")
        _exec_content_box(slide, _EX_LEFT, _EX_COL_BOX_Y,
                          _EX_COL_W, _EX_COL_BOX_H, done_items or [])
        _exec_content_box(slide, _EX_RIGHT_X, _EX_COL_BOX_Y,
                          _EX_COL_W, _EX_COL_BOX_H, next_items or [])

        # ── Zone 4: Attention (fixed position) ──
        _exec_fieldset_label(slide, _EX_LEFT, _EX_ATT_LABEL_Y,
                             "Pontos de Atenção", color=C.AMBER)
        _exec_content_box(slide, _EX_LEFT, _EX_ATT_BOX_Y,
                          _EX_WIDTH, _EX_ATT_BOX_H, risk_items or [])

        # ── Footer (fixed position) ──
        if footnote:
            _textbox(slide, MARGIN, _EX_FOOT_Y, FOOTNOTE_W, 274320,
                     footnote, size=Pt(7), color=C.TEXT_LIGHT)

        # ── Legend (fixed position) ──
        _exec_legend(slide, _EX_LEGEND_Y)

        return slide

    # ── Closing (DARK) ──

    def add_closing(self, title="OBRIGADO", cta=None,
                    footer="M7 Investimentos"):
        """Closing slide with dark background and optional CTA."""
        slide = self._dark_slide()
        _textbox(slide, 914400, 1463040, 7315200, 548640,
                 title, font=FONT_HEADING, size=Pt(24), bold=False,
                 color=C.OFF_WHITE, align=PP_ALIGN.CENTER)
        if cta:
            _textbox(slide, 914400, 2103120, 7315200, 914400,
                     cta, font=FONT_HEADING, size=Pt(20), bold=False,
                     color=C.OFF_WHITE, align=PP_ALIGN.CENTER)
        _separator(slide, 3840480, 3749040, 1463040, C.LIME)
        _textbox(slide, 1828800, 4023360, 5486400, 365760,
                 footer, size=Pt(12), color=C.OFF_WHITE,
                 align=PP_ALIGN.CENTER)
        return slide

    # ── Blank Slide (for full creative freedom) ──

    def add_blank_slide(self, light=True, section_num=None,
                        section_label=None, title=None):
        """Blank slide for custom layouts. Returns SlideBuilder."""
        slide = self._light_slide() if light else self._dark_slide()
        txt_color = C.TEXT_TITLE if light else C.OFF_WHITE
        lbl_color = C.LIME if light else C.OFF_WHITE

        if section_num and section_label:
            _textbox(slide, MARGIN, HEADER_TOP, 5486400, 274320,
                     f"{section_num} — {section_label}",
                     size=Pt(9), bold=True, color=lbl_color)
        if title:
            _textbox(slide, MARGIN, TITLE_TOP, CONTENT_W, 548640,
                     title, font=FONT_HEADING, size=Pt(24), bold=False,
                     color=txt_color)

        sb = SlideBuilder(slide)
        if title:
            sb.y = INTRO_TOP
        elif section_num:
            sb.y = TITLE_TOP
        else:
            sb.y = HEADER_TOP
        return sb

    # ── Save ──

    def save(self, path):
        """Save the presentation."""
        self._prs.save(str(path))
        n = len(self._prs.slides)
        print(f"✓ Saved: {path} ({n} slides)")
        return path
